"""Slide 4 v2 — fix overflow by reducing font to 12pt, extending box heights,
enabling normAutofit, and tightening text. Input: v6 (v5 + slide 4 v1).

Geometry adjustments (in inches):
  Solutions:        y 1.03  -> 0.85,  h 1.53 -> 1.90
  Use Cases:        y 2.98  -> 2.80,  h 1.55 -> 1.85
  Supporting Facts: y 4.82  -> 4.70,  h 2.19 -> 2.35
  Feasibility:      y 0.80  -> 0.75,  h 2.22 -> 2.20  (essentially unchanged)
  Viability:        y 3.24  -> 3.05,  h 2.15 -> 2.15
  Challenges:       y 5.65  -> 5.30,  h 1.49 -> 1.75

  Bottom bar is at y=7.15; all boxes safely end above 7.10.
"""
import zipfile, shutil, os, tempfile
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Inches
from copy import deepcopy

SRC = r"C:\Users\bari2\Desktop\SIH2026\SIH2025-IDEA-Presentation-Format-v6.pptx"
DST = r"C:\Users\bari2\Desktop\SIH2026\SIH2025-IDEA-Presentation-Format-v10.pptx"

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
A = f"{{{A_NS}}}"
P = f"{{{P_NS}}}"

# Each panel: (shape_id, paragraphs [(bold, regular)], font_pt, (top_in, height_in))
PANELS = {
    4: (   # Feasibility
        [
            ("⚖️ Feasibility", ""),
            ("Technical: ",
             "FastAPI + Next.js 16 + ONNX — open-source, no custom hardware."),
            ("Modular: ",
             "Backend, frontend, edge node deploy independently."),
            ("Market: ",
             "DGMS mandates slope monitoring for opencast mines >100 m."),
            ("Economic: ",
             "$200 edge node vs. $250K–$500K SSR; XGBoost live in production."),
        ],
        12.0,
        (0.75, 2.25),
    ),
    5: (   # Viability
        [
            ("💰 Viability & Business Potential", ""),
            ("Production Champion: ",
             "XGBoost: F1 0.985, 0 / 197 missed evacuations."),
            ("Cost Efficiency: ",
             "Sub-$1K per pit vs. $250K–$500K SSR — extends safety perimeter."),
            ("Compliance: ",
             "Dedup'd, auditable alert logs map to DGMS inspections."),
            ("Resilience: ",
             "Edge-capable (Pi 4/5 + ONNX) — never depends on internet."),
        ],
        12.0,
        (3.00, 2.30),
    ),
    14: (  # Challenges
        [
            ("🚧 Challenges", ""),
            ("Synthetic Stream: ",
             "Demo uses Fukuzono-calibrated synthetic data; real IoT next."),
            ("Edge Hardware: ",
             "ONNX scripts designed; Pi procurement and field test pending."),
            ("Multi-Channel Dispatch: ",
             "SMS / WhatsApp architecture ready; Twilio provisioning pending."),
        ],
        12.0,
        (5.35, 1.75),
    ),
    9: (   # Solutions
        [
            ("🛡️ Solutions", ""),
            ("Edge: ",
             "ONNX-exported XGBoost on Raspberry Pi 4/5 — no internet, no cloud."),
            ("Physics: ",
             "Fukuzono (1985) inverse-velocity — same physics as real SSR."),
            ("Training: ",
             "Class-weighted loss → 100% evacuation recall (0 / 197 missed)."),
        ],
        12.0,
        (0.85, 1.90),
    ),
    11: (  # Use Cases
        [
            ("👨‍💻 Use Cases", ""),
            ("Safety Officers: ",
             "Real-time 3D pit heatmap, color-coded risk zones."),
            ("Shift In-Charge: ",
             "Sub-second WebSocket alerts enable truck rerouting."),
            ("Pit Workers: ",
             "Siren + SMS + WhatsApp + push, reach the right person fast."),
        ],
        12.0,
        (2.80, 1.85),
    ),
}

SHAPE_25_PANEL = [
    ("⭐  Supporting Facts for Feasibility and Viability ⭐", ""),
    ("→ ",
     "DGMS mandates slope monitoring for opencast mines >100 m deep — binding regulatory pull."),
    ("→ ",
     "Target: Kusmunda Opencast Mine, SECL, Chhattisgarh — real operating mine with SSR."),
    ("→ ",
     "XGBoost: F1 0.985, 0 / 197 missed evacuations; SHAP terrain 17.03%, SAR 6.90%."),
]
SHAPE_25_GEOM = (4.70, 2.35)  # (top_in, height_in)
SHAPE_25_FONT = 12.0


def find_shape(slide, shape_id):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise SystemExit(f"Shape id {shape_id} not found")


def clear_paragraph(p):
    for r in list(p.runs):
        r._r.getparent().remove(r._r)


def add_bold_then_text(paragraph, bold_part, regular_part, font_size_pt):
    run_b = paragraph.add_run()
    run_b.text = bold_part
    run_b.font.bold = True
    run_b.font.size = Pt(font_size_pt)
    if regular_part:
        run_r = paragraph.add_run()
        run_r.text = regular_part
        run_r.font.bold = False
        run_r.font.size = Pt(font_size_pt)


def enable_normAutofit(txBody, font_scale_pct=100000, lnSpc_reduction_pct=0):
    """Add <a:normAutofit/> to bodyPr so PowerPoint shrinks text to fit."""
    bodyPr = txBody.find(A + "bodyPr")
    if bodyPr is None:
        bodyPr = etree.SubElement(txBody, A + "bodyPr")
    # Remove any existing autofit
    for child in list(bodyPr):
        if etree.QName(child).localname in ("normAutofit", "spAutoFit", "noAutofit"):
            bodyPr.remove(child)
    autofit = etree.SubElement(bodyPr, A + "normAutofit")
    autofit.set("fontScale", str(font_scale_pct))
    if lnSpc_reduction_pct:
        autofit.set("lnSpcReduction", str(lnSpc_reduction_pct))


def rewrite_textbox(shape, paragraphs, font_size_pt, top_in, height_in):
    # Adjust box geometry
    shape.top = Inches(top_in)
    shape.height = Inches(height_in)

    tf = shape.text_frame
    txBody = tf._txBody
    p_elems = txBody.findall(A + "p")
    for extra_p in p_elems[1:]:
        txBody.remove(extra_p)
    seed_para = tf.paragraphs[0]
    clear_paragraph(seed_para)
    add_bold_then_text(seed_para, *paragraphs[0], font_size_pt)
    from pptx.text.text import _Paragraph
    for bold_part, regular_part in paragraphs[1:]:
        new_p = deepcopy(seed_para._p)
        for r in list(new_p.findall(A + "r")):
            new_p.remove(r)
        txBody.append(new_p)
        para = _Paragraph(new_p, tf)
        add_bold_then_text(para, bold_part, regular_part, font_size_pt)
    # Enable normAutofit
    enable_normAutofit(txBody, font_scale_pct=100000, lnSpc_reduction_pct=0)


def add_bold_then_text_xml(parent_p, bold_part, regular_part, font_size_pt):
    run_b = etree.SubElement(parent_p, A + "r")
    rPr = etree.SubElement(run_b, A + "rPr")
    rPr.set("lang", "en-US")
    rPr.set("b", "1")
    rPr.set("sz", str(int(font_size_pt * 100)))
    t_b = etree.SubElement(run_b, A + "t")
    t_b.text = bold_part
    if regular_part:
        run_r = etree.SubElement(parent_p, A + "r")
        rPr_r = etree.SubElement(run_r, A + "rPr")
        rPr_r.set("lang", "en-US")
        rPr_r.set("sz", str(int(font_size_pt * 100)))
        t_r = etree.SubElement(run_r, A + "t")
        t_r.text = regular_part


def fix_shape25_in_xml(pptx_path: str, paragraphs, font_size_pt, top_in, height_in):
    """Patch shape 25 (mc:Choice) and its xfrm via direct XML on the zip."""
    tmp_dir = tempfile.mkdtemp(prefix="pptx_edit_")
    try:
        with zipfile.ZipFile(pptx_path, "r") as zin:
            zin.extractall(tmp_dir)
        slide4_path = os.path.join(tmp_dir, "ppt", "slides", "slide4.xml")
        parser = etree.XMLParser(remove_blank_text=False)
        tree = etree.parse(slide4_path, parser)
        root = tree.getroot()

        # Find the <p:sp> with cNvPr id=25 inside <mc:Choice>
        target_sp = None
        for sp in root.iter(P + "sp"):
            cNvPr = sp.find(P + "nvSpPr/" + P + "cNvPr")
            if cNvPr is None or cNvPr.get("id") != "25":
                continue
            parent = sp.getparent()
            in_choice = False
            while parent is not None:
                qn = etree.QName(parent)
                if qn.namespace == MC_NS:
                    if qn.localname == "Choice":
                        in_choice = True
                        break
                    elif qn.localname == "Fallback":
                        break
                parent = parent.getparent()
            if in_choice:
                target_sp = sp
                break
        if target_sp is None:
            raise SystemExit("shape 25 (mc:Choice) not found")

        # Update geometry (spPr/xfrm/off and ext)
        spPr = target_sp.find(P + "spPr")
        if spPr is not None:
            xfrm = spPr.find(A + "xfrm")
            if xfrm is None:
                xfrm = etree.SubElement(spPr, A + "xfrm")
            off = xfrm.find(A + "off")
            if off is None:
                off = etree.SubElement(xfrm, A + "off")
            off.set("x", str(int(top_in * 914400)))
            off.set("y", str(int(4.30 * 914400)))  # 4.30 is preserved from original y; if we want new y, set here
            # Actually we want to use the SHAPE_25_GEOM
            off.set("y", str(int(4.30 * 914400)))  # placeholder; will be replaced below
        # Better: use SHAPE_25_GEOM values directly
        off = xfrm.find(A + "off")
        off.set("x", str(int(6.98 * 914400)))
        off.set("y", str(int(top_in * 914400)))
        ext = xfrm.find(A + "ext")
        if ext is None:
            ext = etree.SubElement(xfrm, A + "ext")
        ext.set("cx", str(int(6.11 * 914400)))
        ext.set("cy", str(int(height_in * 914400)))

        # Find txBody in p: namespace and rewrite
        txBody = target_sp.find(P + "txBody")
        if txBody is None:
            raise SystemExit("p:txBody not found in shape 25")
        for p in list(txBody.findall(A + "p")):
            txBody.remove(p)
        for bold_part, regular_part in paragraphs:
            p = etree.SubElement(txBody, A + "p")
            add_bold_then_text_xml(p, bold_part, regular_part, font_size_pt)
        # Enable normAutofit
        enable_normAutofit(txBody, font_scale_pct=100000, lnSpc_reduction_pct=0)

        tree.write(slide4_path, xml_declaration=True, encoding="UTF-8", standalone=True)
        new_zip = pptx_path + ".tmp"
        with zipfile.ZipFile(new_zip, "w", zipfile.ZIP_DEFLATED) as zout:
            for foldername, _, filenames in os.walk(tmp_dir):
                for filename in filenames:
                    filepath = os.path.join(foldername, filename)
                    arcname = os.path.relpath(filepath, tmp_dir)
                    zout.write(filepath, arcname)
        os.replace(new_zip, pptx_path)
        print(f"  Fixed shape 25 (mc:Choice): top={top_in}\", h={height_in}\"")
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def main():
    prs = Presentation(SRC)
    slide = prs.slides[3]
    for shape_id, (paragraphs, font_pt, (top, height)) in PANELS.items():
        shape = find_shape(slide, shape_id)
        rewrite_textbox(shape, paragraphs, font_pt, top, height)
        print(f"  Wrote shape id={shape_id}: {len(paragraphs)} paragraphs, top={top}\", h={height}\"")
    prs.save(DST)
    print(f"\nSaved (python-pptx): {DST}")
    fix_shape25_in_xml(DST, SHAPE_25_PANEL, SHAPE_25_FONT, *SHAPE_25_GEOM)
    print(f"\nFinal: {DST}")


if __name__ == "__main__":
    main()
