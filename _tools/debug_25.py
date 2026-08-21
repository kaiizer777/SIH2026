"""Debug why shape 25 txBody isn't being found."""
from pptx import Presentation
from lxml import etree

prs = Presentation(r"C:\Users\bari2\Desktop\SIH2026\SIH2025-IDEA-Presentation-Format-v5.pptx")
slide = prs.slides[3]
spTree = slide.shapes._spTree
print(f"spTree tag: {spTree.tag}")

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"

# Walk all txBody
count = 0
for txBody in spTree.iter(f"{{{A_NS}}}txBody"):
    count += 1
    sp = txBody.getparent()
    if sp is None:
        continue
    sp_tag = etree.QName(sp).localname
    if sp_tag != "sp":
        continue
    cNvPr = sp.find(f"{{{P_NS}}}nvSpPr/{{{P_NS}}}cNvPr")
    if cNvPr is None:
        continue
    sid = cNvPr.get("id")
    if sid != "25":
        continue
    # Found a candidate
    # Walk up
    parent = sp.getparent()
    chain = []
    in_choice = False
    while parent is not None:
        qn = etree.QName(parent)
        chain.append(f"{qn.localname}@{qn.namespace[-15:]}")
        if qn.namespace == MC_NS:
            if qn.localname == "Choice":
                in_choice = True
                break
            elif qn.localname == "Fallback":
                break
        parent = parent.getparent()
    print(f"\ntxBody candidate for id=25:")
    print(f"  Ancestors: {' > '.join(chain)}")
    print(f"  in_choice = {in_choice}")
    print(f"  sp has {len(sp.findall(f'{{{A_NS}}}txBody'))} txBody direct child")

print(f"\nTotal txBody found: {count}")
