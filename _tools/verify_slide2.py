from pptx import Presentation
p = Presentation(r'C:\Users\bari2\Desktop\SIH2026\SIH2025-IDEA-Presentation-Format-rewritten.pptx')
for shape in p.slides[1].shapes:
    if shape.has_text_frame and shape.shape_id == 13:
        for i, para in enumerate(shape.text_frame.paragraphs):
            print(f"[{i}]", repr(para.text))
        break
