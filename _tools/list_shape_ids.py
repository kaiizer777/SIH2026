"""List all shape ids that python-pptx sees on slide 4."""
from pptx import Presentation
prs = Presentation(r"C:\Users\bari2\Desktop\SIH2026\SIH2025-IDEA-Presentation-Format-v5.pptx")
slide = prs.slides[3]
for shape in slide.shapes:
    try:
        sid = shape.shape_id
    except Exception as e:
        sid = f"<err: {e}>"
    name = getattr(shape, 'name', '?')
    has_tf = shape.has_text_frame
    text = ''
    if has_tf:
        text = shape.text_frame.text[:80].replace('\n', ' | ')
    print(f"shape_id={sid!s:>8}  name={name!r}  has_tf={has_tf}  text={text!r}")
