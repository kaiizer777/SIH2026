"""Debug: find every txBody whose sp has cNvPr id=25, dump its parent chain."""
from pptx import Presentation
from lxml import etree

prs = Presentation(r"C:\Users\bari2\Desktop\SIH2026\SIH2025-IDEA-Presentation-Format-v5.pptx")
slide = prs.slides[3]
spTree = slide.shapes._spTree

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"

for el in spTree.iter():
    if etree.QName(el).localname != "txBody" or etree.QName(el).namespace != A_NS:
        continue
    sp = el.getparent()
    if sp is None or etree.QName(sp).localname != "sp":
        continue
    cNvPr = sp.find(f"{{{P_NS}}}nvSpPr/{{{P_NS}}}cNvPr")
    if cNvPr is None or cNvPr.get("id") != "25":
        continue
    # Found
    print(f"Found txBody in sp with id=25")
    # Walk up and check
    parent = sp.getparent()
    chain = []
    while parent is not None:
        qn = etree.QName(parent)
        chain.append((qn.localname, qn.namespace))
        if qn.namespace == MC_NS:
            if qn.localname == "Choice":
                print(f"  In Choice! chain={chain}")
                break
            elif qn.localname == "Fallback":
                print(f"  In Fallback! chain={chain}")
                break
        parent = parent.getparent()
    else:
        print(f"  Did NOT hit MC ancestor! Full chain={chain}")
