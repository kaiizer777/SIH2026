"""Rewrite slide 4 of the v5 deck — direct XML on the zip, correct namespaces.

Shape 25's <p:txBody> is in the presentation (p:) namespace, not drawing (a:).
"""
import zipfile, shutil, os, tempfile
from lxml import etree
from pptx import Presentation
from pptx.util import Pt
from copy import deepcopy

SRC = r"C:\Users\bari2\Desktop\SIH2026\SIH2025-IDEA-Presentation-Format-v5.pptx"
DST = r"C:\Users\bari2\Desktop\SIH2026\SIH2025-IDEA-Presentation-Format-v6.pptx"

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
A = f"{{{A_NS}}}"
P = f"{{{P_NS}}}"

PANELS = {
    4: (
        [
            ("⚖️ Feasibility", ""),
            ("Technical: ",
             "Production stack (FastAPI + Next.js 16 + ONNX Runtime) — no custom hardware, all open-source."),
            ("Modular: ",
             "Backend (Render), frontend (Vercel), and edge node deploy independently, scale per pit."),
            ("Market: ",
             "DGMS mandates systematic slope monitoring for opencast mines >100 m deep."),
            ("Economic: ",
             "$200 edge node replaces $250K–$500K SSR units; class-weighted XGBoost runs in production."),
        ],
        None,
    ),
    5: (
        [
            ("💰 Viability & Business Potential", ""),
            ("Production Champion: ",
             "XGBoost delivers F1 0.985 with 0 / 197 missed evacuations on the held-out test set."),
            ("Cost Efficiency: ",
             "Sub-$1K per-pit hardware vs. $250K–$500K SSR — extends the safety perimeter at a fraction of the cost."),
            ("Compliance: ",
             "Dedup'd, auditable alert logs map directly to DGMS inspection and post-incident review."),
            ("Resilience: ",
             "Edge-capable by design (Raspberry Pi 4/5 + ONNX Runtime) — the safety system never depends on internet."),
        ],
        None,
    ),
    14: (
        [
            ("🚧 Challenges", ""),
            ("Synthetic Sensor Stream: ",
             "Live demo uses Fukuzono-calibrated synthetic data; physical IoT integration is the next build step."),
            ("Edge Hardware: ",
             "ONNX conversion scripts designed; Raspberry Pi procurement and field test pending."),
            ("Multi-Channel Dispatch: ",
             "Architecture supports SMS / WhatsApp / siren / push; Twilio gateway provisioning pending."),
        ],
        16.5,
    ),
    9: (
        [
            ("🛡️ Solutions", ""),
            ("Edge Processing: ",
             "ONNX-exported XGBoost runs on Raspberry Pi 4/5 at the bench — no internet, no cloud."),
            ("Physics-Grounded ML: ",
             "Fukuzono (1985) inverse-velocity model — the same physics used in real SSR systems."),
            ("Class-Weighted Training: ",
             "Imbalance-aware loss gives 100% evacuation recall (0 / 197 missed) on the test set."),
        ],
        17.0,
    ),
    11: (
        [
            ("👨‍💻 Use Cases", ""),
            ("Safety Officers: ",
             "Real-time 3D pit heatmap (Next.js 16 + MapLibre GL) with color-coded risk zones."),
            ("Shift In-Charge: ",
             "Sub-second WebSocket alerts enable truck rerouting from at-risk benches."),
            ("Pit Workers: ",
             "Multi-channel alerts (siren + SMS + WhatsApp + push) reach the right person, fast."),
        ],
        None,
    ),
}

SHAPE_25_PANEL = [
    ("⭐  Supporting Facts for Feasibility and Viability ⭐", ""),
    ("→ ",
     "DGMS mandates systematic slope monitoring for all opencast mines >100 m deep — a binding regulatory pull."),
    ("→ ",
     "Live target site: Kusmunda Opencast Mine, SECL, Chhattisgarh — a real operating mine with existing SSR deployment."),
    ("→ ",
     "Production XGBoost: F1 0.985, 0 / 197 missed evacuations; SHAP shows terrain (17.03%) and SAR (6.90%) as dominant signals."),
]


def find_shape(slide, shape_id):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise SystemExit(f"Shape id {shape_id} not found")


def clear_paragraph(p):
    for r in list(p.runs):
        r._r.getparent().remove(r._r)


def add_bold_then_text(paragraph, bold_part, regular_part, font_size_pt):
    run_b = paragraph.add_run()
    run_b.text = bold_part
    run_b.font.bold = True
    if font_size_pt is not None:
        run_b.font.size = Pt(font_size_pt)
    if regular_part:
        run_r = paragraph.add_run()
        run_r.text = regular_part
        run_r.font.bold = False
        if font_size_pt is not None:
            run_r.font.size = Pt(font_size_pt)


def rewrite_textbox(shape, paragraphs, font_size_pt):
    tf = shape.text_frame
    txBody = tf._txBody
    p_elems = txBody.findall(A + "p")
    for extra_p in p_elems[1:]:
        txBody.remove(extra_p)
    seed_para = tf.paragraphs[0]
    clear_paragraph(seed_para)
    add_bold_then_text(seed_para, *paragraphs[0], font_size_pt)
    from pptx.text.text import _Paragraph
    for bold_part, regular_part in paragraphs[1:]:
        new_p = deepcopy(seed_para._p)
        for r in list(new_p.findall(A + "r")):
            new_p.remove(r)
        txBody.append(new_p)
        para = _Paragraph(new_p, tf)
        add_bold_then_text(para, bold_part, regular_part, font_size_pt)


def add_bold_then_text_xml(parent_p, bold_part, regular_part, font_size_pt):
    """Add bold + optional regular run as lxml elements under a parent <a:p>."""
    run_b = etree.SubElement(parent_p, A + "r")
    rPr = etree.SubElement(run_b, A + "rPr")
    rPr.set("lang", "en-US")
    rPr.set("b", "1")
    if font_size_pt is not None:
        rPr.set("sz", str(int(font_size_pt * 100)))
    t_b = etree.SubElement(run_b, A + "t")
    t_b.text = bold_part
    if regular_part:
        run_r = etree.SubElement(parent_p, A + "r")
        rPr_r = etree.SubElement(run_r, A + "rPr")
        rPr_r.set("lang", "en-US")
        if font_size_pt is not None:
            rPr_r.set("sz", str(int(font_size_pt * 100)))
        t_r = etree.SubElement(run_r, A + "t")
        t_r.text = regular_part


def fix_shape25_in_xml(pptx_path: str):
    """Find the <p:txBody> in shape 25 (mc:Choice variant) and rewrite it."""
    tmp_dir = tempfile.mkdtemp(prefix="pptx_edit_")
    try:
        with zipfile.ZipFile(pptx_path, "r") as zin:
            zin.extractall(tmp_dir)
        slide4_path = os.path.join(tmp_dir, "ppt", "slides", "slide4.xml")
        parser = etree.XMLParser(remove_blank_text=False)
        tree = etree.parse(slide4_path, parser)
        root = tree.getroot()

        # txBody in shape 25 inside <mc:Choice> is in the p: (presentation) namespace
        target_txBody = None
        for txBody in root.iter(P + "txBody"):
            sp = txBody.getparent()
            if sp is None or etree.QName(sp).localname != "sp":
                continue
            cNvPr = sp.find(P + "nvSpPr/" + P + "cNvPr")
            if cNvPr is None or cNvPr.get("id") != "25":
                continue
            parent = sp.getparent()
            in_choice = False
            while parent is not None:
                qn = etree.QName(parent)
                if qn.namespace == MC_NS:
                    if qn.localname == "Choice":
                        in_choice = True
                        break
                    elif qn.localname == "Fallback":
                        break
                parent = parent.getparent()
            if in_choice:
                target_txBody = txBody
                break

        if target_txBody is None:
            raise SystemExit("shape 25 (mc:Choice) <p:txBody> not found")

        for p in list(target_txBody.findall(A + "p")):
            target_txBody.remove(p)
        for bold_part, regular_part in SHAPE_25_PANEL:
            p = etree.SubElement(target_txBody, A + "p")
            add_bold_then_text_xml(p, bold_part, regular_part, 16.0)

        tree.write(slide4_path, xml_declaration=True, encoding="UTF-8", standalone=True)

        new_zip = pptx_path + ".tmp"
        with zipfile.ZipFile(new_zip, "w", zipfile.ZIP_DEFLATED) as zout:
            for foldername, _, filenames in os.walk(tmp_dir):
                for filename in filenames:
                    filepath = os.path.join(foldername, filename)
                    arcname = os.path.relpath(filepath, tmp_dir)
                    zout.write(filepath, arcname)
        os.replace(new_zip, pptx_path)
        print("  Fixed shape 25 (mc:Choice) in slide4.xml directly")
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def main():
    prs = Presentation(SRC)
    slide = prs.slides[3]
    for shape_id, (paragraphs, font_pt) in PANELS.items():
        shape = find_shape(slide, shape_id)
        rewrite_textbox(shape, paragraphs, font_pt)
        print(f"  Wrote shape id={shape_id}: {len(paragraphs)} paragraphs")
    prs.save(DST)
    print(f"\nSaved (python-pptx): {DST}")
    fix_shape25_in_xml(DST)
    print(f"\nFinal: {DST}")


if __name__ == "__main__":
    main()
