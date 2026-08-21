"""Debug: print cNvPr ids of all sps that have txBody."""
from pptx import Presentation
from lxml import etree

prs = Presentation(r"C:\Users\bari2\Desktop\SIH2026\SIH2025-IDEA-Presentation-Format-v5.pptx")
slide = prs.slides[3]
spTree = slide.shapes._spTree

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

print("All sps with txBody (id, name):")
for el in spTree.iter():
    if etree.QName(el).localname != "sp":
        continue
    cNvPr = el.find(f"{{{P_NS}}}nvSpPr/{{{P_NS}}}cNvPr")
    if cNvPr is None:
        continue
    sid = cNvPr.get("id")
    name = cNvPr.get("name")
    has_tx = el.find(f"{{{A_NS}}}txBody") is not None
    if has_tx:
        print(f"  id={sid!r}  name={name!r}")
    # Also check descendants
    for sub_tx in el.iter(f"{{{A_NS}}}txBody"):
        print(f"    -> sub txBody found inside id={sid!r}")
        break
