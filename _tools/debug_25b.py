"""Debug: count elements at each depth."""
from pptx import Presentation
from lxml import etree

prs = Presentation(r"C:\Users\bari2\Desktop\SIH2026\SIH2025-IDEA-Presentation-Format-v5.pptx")
slide = prs.slides[3]
spTree = slide.shapes._spTree

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"

# Count by tag
tags = {}
for el in spTree.iter():
    tag = etree.QName(el).localname
    tags[tag] = tags.get(tag, 0) + 1
for tag, n in sorted(tags.items(), key=lambda x: -x[1]):
    print(f"  {tag}: {n}")
