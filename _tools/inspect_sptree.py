"""Inspect python-pptx's _spTree object to find the lxml root."""
from pptx import Presentation
prs = Presentation(r"C:\Users\bari2\Desktop\SIH2026\SIH2025-IDEA-Presentation-Format-v5.pptx")
slide = prs.slides[3]
spTree = slide.shapes._spTree
print(f"type: {type(spTree)}")
print(f"class: {spTree.__class__.__name__}")
print(f"class mro: {[c.__name__ for c in type(spTree).__mro__]}")
print(f"has _element: {hasattr(spTree, '_element')}")
print(f"has element: {hasattr(spTree, 'element')}")
print(f"is lxml ElementBase: {isinstance(spTree, type(spTree).__mro__[-2])}")
# Check if it has iter directly
print(f"has iter: {hasattr(spTree, 'iter')}")
# Try to find the slide part
print(f"part: {spTree.part if hasattr(spTree, 'part') else 'n/a'}")
