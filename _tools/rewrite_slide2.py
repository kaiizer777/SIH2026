"""Rewrite slide 2 body of the SIH PPT (v5) with overflow fix.

Fixes applied:
  1. Body font reduced 14pt -> 12pt (still readable in left half of LAYOUT_WIDE).
  2. Line spacing tightened to 1.0 (was ~1.2 default).
  3. Body box height extended into empty space (top 1.10, bottom 6.85, h=5.75").
  4. normAutofit enabled so PowerPoint shrinks text further if still over.
  5. Hard-compressed Pillar 3 and Pillar 4 to keep the bullets punchy.
"""
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from copy import deepcopy
from lxml import etree

SRC = r"C:\Users\bari2\.minimax\v2\assets\2026\08\20\21-47-41-438-asset_20260820-214741-438_68aa38d662ed_5ea633b1-SIH2025-IDEA-Presentation-Format.pptx"
DST = r"C:\Users\bari2\Desktop\SIH2026\SIH2025-IDEA-Presentation-Format-v6.pptx"

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
A = f"{{{A_NS}}}"

EM = "\u2014"   # em dash
EN = "\u2013"   # en dash
MD = "\u00b7"   # middle dot
BODY_FONT_PT = 12   # was 14
LINE_SPACING = 1.0  # tighter than 1.2 default

# Hard-compressed body text
SLIDE2_BODY = [
    # --- Lead ---
    ("The system. ",
     "AlertX " + EM + " an end-to-end platform fusing Earth-observation data, physics-grounded sensor telemetry, and class-weighted ML to predict slope failure and dispatch sub-minute alerts."),

    # --- Pillar 1 ---
    ("Pillar 1 " + MD + " Multi-Modal Data Ingestion. ",
     "Sentinel-1 SAR, Copernicus GLO-30 DEM (slope, aspect, curvature), Open-Meteo ERA5 rainfall, and Fukuzono-calibrated synthetic sensors " + EM + " live over FastAPI + WebSocket."),

    # --- Pillar 2 ---
    ("Pillar 2 " + MD + " Physics-Informed ML Core. ",
     "XGBoost (production champion: 0 / 197 missed evacuations, F1 0.985), RandomForest (stronger terrain/SAR SHAP), GRU benchmark " + EM + " all ONNX-exportable for offline edge."),

    # --- Pillar 3 ---
    ("Pillar 3 " + MD + " Risk Classification + Dispatch. ",
     "Three-tier Safe / Warning / Evacuation (peer-reviewed SSR field study) feed a Next.js 16 + MapLibre GL heatmap and SMS / WhatsApp / siren / push channels with sub-second latency."),

    # --- Pillar 4 ---
    ("Pillar 4 " + MD + " Edge-Ready by Design. ",
     "Sensor to inference to siren runs offline on a $" + "200 Raspberry Pi 4/5 at the bench. ONNX Runtime (over TFLite for ARM). Connectivity is for the dashboard, never for safety."),

    # --- Closer ---
    ("We don't replace SSR " + EM + " we extend the safety perimeter to every bench and ramp.", ""),
]


def find_body_shape(slide, shape_id: int):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise SystemExit(f"Shape id {shape_id} not found on slide 2")


def clear_paragraph(p):
    for r in list(p.runs):
        r._r.getparent().remove(r._r)


def add_bold_then_text(paragraph, bold_part: str, regular_part: str, font_size_pt: int = BODY_FONT_PT):
    run_b = paragraph.add_run()
    run_b.text = bold_part
    run_b.font.bold = True
    run_b.font.size = Pt(font_size_pt)
    if regular_part:
        run_r = paragraph.add_run()
        run_r.text = regular_part
        run_r.font.bold = False
        run_r.font.size = Pt(font_size_pt)


def main():
    prs = Presentation(SRC)
    slide = prs.slides[1]
    body = find_body_shape(slide, 13)
    tf = body.text_frame

    # --- 1. Extend the box: top=1.10, height=5.75, width keeps original 6.45 ---
    body.top = Inches(1.10)
    body.height = Inches(5.75)
    # body.left and body.width are unchanged

    # --- 2. Set tighter line spacing on every paragraph (existing + new) ---
    # We'll set this on the txBody default paragraph properties (lstStyle) and on each <a:p>.
    txBody = tf._txBody
    # Remove any existing lstStyle / normAutofit we'll replace
    for child in list(txBody):
        tag = etree.QName(child).localname
        if tag in ("lstStyle",):
            txBody.remove(child)

    # Insert a fresh lstStyle right after bodyPr with our line spacing.
    # But bodyPr comes first in txBody. We'll add lstStyle as the 2nd child if bodyPr exists, else first.
    bodyPr = txBody.find(A + "bodyPr")
    lstStyle = etree.SubElement(txBody, A + "lstStyle") if bodyPr is None else None
    if lstStyle is None:
        lstStyle = etree.Element(A + "lstStyle")
        bodyPr.addnext(lstStyle)
    # default paragraph props
    defPpr = etree.SubElement(lstStyle, A + "defPPr")
    lnSpc = etree.SubElement(defPpr, A + "lnSpc")
    spcPct = etree.SubElement(lnSpc, A + "spcPct")
    spcPct.set("val", str(int(LINE_SPACING * 100000)))  # 100000 = 100% = 1.0

    # --- 3. Wipe existing paragraphs and rebuild ---
    first_p = tf.paragraphs[0]
    p_elems = txBody.findall(A + "p")
    for extra_p in p_elems[1:]:
        txBody.remove(extra_p)
    clear_paragraph(first_p)
    add_bold_then_text(first_p, *SLIDE2_BODY[0])

    for lead, tail in SLIDE2_BODY[1:]:
        new_p = deepcopy(first_p._p)
        for r in list(new_p.findall(A + "r")):
            new_p.remove(r)
        txBody.append(new_p)
        from pptx.text.text import _Paragraph
        para = _Paragraph(new_p, tf)
        add_bold_then_text(para, lead, tail)

    # --- 4. Enable normAutofit (shrink text on overflow) on bodyPr ---
    if bodyPr is not None:
        # Remove any existing autofit children
        for child in list(bodyPr):
            tag = etree.QName(child).localname
            if tag in ("normAutofit", "spAutoFit", "noAutofit"):
                bodyPr.remove(child)
        normAutofit = etree.SubElement(bodyPr, A + "normAutofit")
        normAutofit.set("fontScale", "92500")   # start at 92.5%; PowerPoint will adjust
        normAutofit.set("lnSpcReduction", "10000")  # 10% line-space reduction

    prs.save(DST)
    print(f"Wrote: {DST}")


if __name__ == "__main__":
    main()
